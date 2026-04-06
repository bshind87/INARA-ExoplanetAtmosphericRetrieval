#!/usr/bin/env python
"""
Build the INARA project presentation as a .pptx file.
Run:  python docs/build_presentation.py
Output: docs/INARA_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import copy
from pathlib import Path

# ── Colour palette ────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x00, 0x2B, 0x5C)   # Northeastern navy
RED       = RGBColor(0xCC, 0x00, 0x00)   # Northeastern red
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF0, 0xF4, 0xF9)  # very light blue
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
MID_GREY  = RGBColor(0x55, 0x55, 0x55)
ACCENT1   = RGBColor(0x1A, 0x6B, 0xA0)  # steel blue (baseline)
ACCENT2   = RGBColor(0xB2, 0x18, 0x2B)  # crimson (deep model)
ACCENT3   = RGBColor(0x1A, 0x7A, 0x3C)  # forest green
LIGHT_NAVY = RGBColor(0xD6, 0xE4, 0xF0)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Helpers ───────────────────────────────────────────────────────────────────
def add_rect(slide, left, top, width, height, fill_color, transparency=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if transparency:
        shape.fill.fore_color.theme_color  # access to enable
    shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=DARK_TEXT,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_para(tf, text, font_size=16, bold=False, color=DARK_TEXT,
             align=PP_ALIGN.LEFT, space_before=0, italic=False, level=0):
    from pptx.util import Pt as _Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.level = level
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = _Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def slide_header(slide, title, subtitle=None, title_size=28):
    """Navy top bar with white title."""
    add_rect(slide, 0, 0, 13.33, 1.1, NAVY)
    add_textbox(slide, title, 0.35, 0.12, 12.0, 0.7,
                font_size=title_size, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, subtitle, 0.35, 0.78, 10.0, 0.32,
                    font_size=13, bold=False, color=RGBColor(0xCC, 0xDD, 0xFF),
                    align=PP_ALIGN.LEFT)
    # Red accent bar below header
    add_rect(slide, 0, 1.1, 13.33, 0.06, RED)


def footer(slide, text="CS 6140 · ML · Northeastern University · Spring 2026"):
    add_rect(slide, 0, 7.15, 13.33, 0.35, NAVY)
    add_textbox(slide, text, 0.3, 7.16, 12.0, 0.28,
                font_size=9, color=RGBColor(0xAA, 0xBB, 0xDD), align=PP_ALIGN.LEFT)


def bullet_box(slide, items, left, top, width, height,
               font_size=15, bullet="▸ ", color=DARK_TEXT, line_spacing=6):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        if line_spacing and not first:
            p.space_before = Pt(line_spacing)
        run = p.add_run()
        if isinstance(item, tuple):
            run.text = item[0]
            run.font.size = Pt(item[1])
            run.font.bold = item[2] if len(item) > 2 else False
            run.font.color.rgb = item[3] if len(item) > 3 else color
        else:
            run.text = bullet + item
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
    return txb


def card(slide, left, top, width, height, fill=LIGHT_NAVY,
         title=None, title_color=NAVY, body_items=None, body_size=13):
    add_rect(slide, left, top, width, height, fill)
    y = top + 0.12
    if title:
        add_textbox(slide, title, left+0.15, y, width-0.3, 0.38,
                    font_size=14, bold=True, color=title_color)
        y += 0.38
    if body_items:
        for item in body_items:
            add_textbox(slide, "  " + item, left+0.1, y, width-0.2, 0.32,
                        font_size=body_size, color=DARK_TEXT)
            y += 0.3


def two_col_table(slide, headers, rows, left, top, width, height,
                  col_widths=None, header_fill=NAVY, row_fill=LIGHT_BG,
                  alt_fill=WHITE, font_size=12):
    n_cols = len(headers)
    if col_widths is None:
        col_widths = [width / n_cols] * n_cols
    row_h = height / (len(rows) + 1)

    # Header row
    x = left
    for i, (h, cw) in enumerate(zip(headers, col_widths)):
        add_rect(slide, x, top, cw, row_h, header_fill)
        add_textbox(slide, h, x+0.05, top+0.04, cw-0.1, row_h-0.05,
                    font_size=font_size, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x += cw

    # Data rows
    for ri, row in enumerate(rows):
        fill = row_fill if ri % 2 == 0 else alt_fill
        x = left
        for ci, (val, cw) in enumerate(zip(row, col_widths)):
            add_rect(slide, x, top + row_h * (ri+1), cw, row_h, fill)
            add_textbox(slide, str(val),
                        x+0.05, top + row_h*(ri+1) + 0.04,
                        cw-0.1, row_h-0.05,
                        font_size=font_size-1, color=DARK_TEXT,
                        align=PP_ALIGN.CENTER)
            x += cw


# ── Build presentation ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # blank


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)

# Background gradient effect — two rects
add_rect(sl, 0, 0, 13.33, 7.5, NAVY)
add_rect(sl, 0, 4.8, 13.33, 2.7, RGBColor(0x00, 0x1A, 0x40))

# Red accent stripe
add_rect(sl, 0, 2.0, 0.12, 2.8, RED)

# University + course
add_textbox(sl, "NORTHEASTERN UNIVERSITY  ·  CS 6140 MACHINE LEARNING  ·  SPRING 2026",
            0.3, 0.28, 12.5, 0.4, font_size=11, color=RGBColor(0xAA, 0xBB, 0xDD),
            bold=True)

# Title
add_textbox(sl, "Exploring ML Models for Detection of",
            0.3, 1.0, 12.5, 0.75, font_size=34, bold=True, color=WHITE)
add_textbox(sl, "Atmospheric Composition in Exoplanets",
            0.3, 1.7, 12.5, 0.75, font_size=34, bold=True, color=RGBColor(0xFF, 0xDD, 0x88))

# Subtitle
add_textbox(sl, "Machine learning–based atmospheric retrieval from synthetic transmission spectra",
            0.3, 2.7, 12.0, 0.5, font_size=16, italic=True,
            color=RGBColor(0xCC, 0xDD, 0xFF))

# Team
add_textbox(sl, "Team Members", 0.3, 3.55, 4.0, 0.35,
            font_size=13, bold=True, color=RGBColor(0xFF, 0xDD, 0x88))

team = [
    "Shantanu Wankhare   wankhare.s@northeastern.edu",
    "Bhalchandra Shinde   shinde.b@northeastern.edu",
    "Asad Mulani             mulani.a@northeastern.edu",
]
for i, t in enumerate(team):
    add_textbox(sl, t, 0.3, 3.95 + i*0.38, 7.0, 0.36,
                font_size=14, color=WHITE)

# Dataset pill
add_rect(sl, 0.3, 5.45, 3.8, 0.52, ACCENT1)
add_textbox(sl, "Dataset: INARA  (NASA FDL / Zorzan et al. 2025)",
            0.35, 5.48, 3.7, 0.42, font_size=13, bold=True, color=WHITE)

# Star/telescope icon placeholder
add_rect(sl, 9.8, 3.3, 3.2, 3.2, RGBColor(0x00, 0x1F, 0x4A))
add_textbox(sl, "🔭", 10.4, 3.7, 2.0, 2.0, font_size=72, color=WHITE)

footer(sl)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — MOTIVATION
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, WHITE)
slide_header(sl, "Motivation — Why This Matters",
             "The intersection of scientific urgency and methodological opportunity")
footer(sl)

# Left column — scientific urgency
add_rect(sl, 0.3, 1.35, 5.8, 5.5, LIGHT_BG)
add_textbox(sl, "The Science", 0.5, 1.45, 5.4, 0.4,
            font_size=16, bold=True, color=NAVY)

sci = [
    "Rocky exoplanet atmospheres reveal if a world could be habitable",
    "Co-detection of O₂ + CH₄ is a strong biosignature — signals biological activity",
    "JWST + planned Habitable Worlds Observatory (HWO) produce spectra at unprecedented scale",
    "Population-level surveys require analysing thousands of planets simultaneously",
]
for i, s in enumerate(sci):
    add_textbox(sl, "▸  " + s, 0.45, 1.92 + i*0.88, 5.5, 0.82,
                font_size=13.5, color=DARK_TEXT)

# Right column — the problem with traditional methods
add_rect(sl, 6.5, 1.35, 6.5, 2.5, RGBColor(0xFF, 0xEE, 0xEE))
add_textbox(sl, "Traditional Bayesian Retrieval", 6.7, 1.45, 6.1, 0.4,
            font_size=16, bold=True, color=ACCENT2)
trad = [
    "Nested sampling / MCMC — rigorous but slow",
    "Hours to days per planet",
    "Does not scale to thousands of targets",
]
for i, t in enumerate(trad):
    add_textbox(sl, "✗  " + t, 6.65, 1.92 + i*0.58, 6.1, 0.52,
                font_size=13.5, color=DARK_TEXT)

# Right column — ML opportunity
add_rect(sl, 6.5, 4.1, 6.5, 2.6, RGBColor(0xE8, 0xF5, 0xE9))
add_textbox(sl, "ML-Based Retrieval", 6.7, 4.2, 6.1, 0.4,
            font_size=16, bold=True, color=ACCENT3)
ml = [
    "Milliseconds per planet once trained",
    "Enables population-level atmospheric studies",
    "Active research direction (JWST era, 2022–present)",
    "Established benchmarks exist (Márquez-Neila 2018; Vasist 2023; Gebhard 2024)",
]
for i, t in enumerate(ml):
    add_textbox(sl, "✓  " + t, 6.65, 4.68 + i*0.50, 6.1, 0.46,
                font_size=13.5, color=DARK_TEXT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, WHITE)
slide_header(sl, "Problem Statement", "Multi-output regression over molecular log₁₀ abundances")
footer(sl)

# Central problem box
add_rect(sl, 1.0, 1.35, 11.33, 1.4, LIGHT_NAVY)
add_textbox(sl, "Given a transmission spectrum of an exoplanet's atmosphere,",
            1.2, 1.42, 10.9, 0.45, font_size=17, bold=False, color=NAVY)
add_textbox(sl, "predict the log₁₀ surface volume mixing ratios of 12 molecular species.",
            1.2, 1.85, 10.9, 0.45, font_size=17, bold=True, color=NAVY)

# Input → Model → Output flow
# Input box
add_rect(sl, 0.4, 3.0, 3.8, 2.8, LIGHT_BG)
add_textbox(sl, "INPUT", 0.6, 3.08, 3.4, 0.38,
            font_size=13, bold=True, color=ACCENT1, align=PP_ALIGN.CENTER)
add_textbox(sl, "Transmission Spectrum", 0.55, 3.5, 3.6, 0.38,
            font_size=14, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
inputs = ["3 spectral channels", "4,378 wavelength points", "0.2 – 2.0 µm range", "(HWO / LUVOIR design)"]
for i, inp in enumerate(inputs):
    add_textbox(sl, inp, 0.55, 3.95 + i*0.38, 3.6, 0.35,
                font_size=12, color=MID_GREY, align=PP_ALIGN.CENTER)

# Arrow
add_textbox(sl, "→", 4.35, 4.0, 0.7, 0.6, font_size=36, bold=True, color=NAVY,
            align=PP_ALIGN.CENTER)

# Model box
add_rect(sl, 5.15, 3.0, 3.0, 2.8, NAVY)
add_textbox(sl, "ML MODEL", 5.25, 3.08, 2.8, 0.38,
            font_size=13, bold=True, color=RGBColor(0xFF, 0xDD, 0x88), align=PP_ALIGN.CENTER)
add_textbox(sl, "Random Forest", 5.25, 3.55, 2.8, 0.35,
            font_size=13, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(sl, "(baseline, ≤10k samples)", 5.25, 3.9, 2.8, 0.35,
            font_size=11, italic=True, color=RGBColor(0xAA, 0xBB, 0xFF), align=PP_ALIGN.CENTER)
add_textbox(sl, "1D ResNet", 5.25, 4.3, 2.8, 0.35,
            font_size=13, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(sl, "(deep model, full dataset)", 5.25, 4.65, 2.8, 0.35,
            font_size=11, italic=True, color=RGBColor(0xAA, 0xBB, 0xFF), align=PP_ALIGN.CENTER)

# Arrow
add_textbox(sl, "→", 8.3, 4.0, 0.7, 0.6, font_size=36, bold=True, color=NAVY,
            align=PP_ALIGN.CENTER)

# Output box
add_rect(sl, 9.1, 3.0, 3.9, 2.8, RGBColor(0xE8, 0xF5, 0xE9))
add_textbox(sl, "OUTPUT", 9.3, 3.08, 3.5, 0.38,
            font_size=13, bold=True, color=ACCENT3, align=PP_ALIGN.CENTER)
add_textbox(sl, "12 Molecular Abundances", 9.3, 3.5, 3.5, 0.38,
            font_size=14, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
mols_disp = ["H₂O · CO₂ · O₂ · O₃", "CH₄ · N₂ · N₂O · CO", "H₂ · H₂S · SO₂ · NH₃", "(log₁₀ VMR at surface)"]
for i, m in enumerate(mols_disp):
    add_textbox(sl, m, 9.3, 3.95 + i*0.38, 3.5, 0.35,
                font_size=12, color=MID_GREY, align=PP_ALIGN.CENTER)

# Why log10
add_rect(sl, 0.4, 6.05, 12.5, 0.7, RGBColor(0xFF, 0xF8, 0xE1))
add_textbox(sl, "Why log₁₀ space?  Mixing ratios span ~12 orders of magnitude (~0.1 for N₂ to 10⁻⁴⁰ for trace species). "
            "Log₁₀ makes the prediction scale uniform — an error of 1.0 means 'off by one order of magnitude'.",
            0.6, 6.08, 12.1, 0.62, font_size=12.5, italic=True, color=RGBColor(0x66, 0x44, 0x00))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — DATASET
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, WHITE)
slide_header(sl, "Dataset — INARA (NASA FDL / Zorzan et al. 2025)",
             "3.1 million synthetic rocky planet spectra from NASA Planetary Spectrum Generator")
footer(sl)

# Left: dataset facts
add_rect(sl, 0.3, 1.35, 5.8, 5.6, LIGHT_BG)
add_textbox(sl, "Dataset Facts", 0.5, 1.45, 5.4, 0.38,
            font_size=15, bold=True, color=NAVY)
facts = [
    ("Full dataset size", "3,112,620 synthetic spectra"),
    ("This project uses", "~124,000 samples (limited by compute)"),
    ("Wavelength range", "0.2 – 2.0 µm (HWO / LUVOIR range)"),
    ("Wavelength points", "4,378 per spectrum"),
    ("Stellar types", "F, G, K, M — 4 types"),
    ("Aux features", "8 per system (Tₛ, Rₚ, g, distance, …)"),
    ("Format", "CSV per planet → .npy arrays after processing"),
]
for i, (k, v) in enumerate(facts):
    y = 1.95 + i*0.68
    add_textbox(sl, k, 0.45, y, 2.3, 0.3, font_size=12, bold=True, color=NAVY)
    add_textbox(sl, v, 2.8, y, 3.2, 0.3, font_size=12, color=DARK_TEXT)

# Right top: spectral channels
add_rect(sl, 6.45, 1.35, 6.5, 2.75, LIGHT_NAVY)
add_textbox(sl, "Input Spectral Channels", 6.65, 1.45, 6.1, 0.38,
            font_size=15, bold=True, color=NAVY)
channels = [
    ("Ch 0", "Observed SNR spectrum  (star+planet / noise), normalised per spectrum"),
    ("Ch 1", "Mean-subtracted SNR — removes stellar continuum, isolates molecular absorption"),
    ("Ch 2", "Log transit depth — direct measure of atmosphere thickness"),
]
for i, (ch, desc) in enumerate(channels):
    y = 1.92 + i*0.72
    add_rect(sl, 6.5, y, 0.55, 0.55, ACCENT1)
    add_textbox(sl, ch, 6.52, y+0.08, 0.5, 0.38,
                font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, desc, 7.15, y+0.04, 5.6, 0.5,
                font_size=12, color=DARK_TEXT)

# Right bottom: 12 target molecules
add_rect(sl, 6.45, 4.3, 6.5, 2.55, RGBColor(0xE8, 0xF5, 0xE9))
add_textbox(sl, "12 Target Molecules", 6.65, 4.4, 6.1, 0.38,
            font_size=15, bold=True, color=ACCENT3)
mol_groups = [
    ("Habitability", "H₂O · CO₂", ACCENT1),
    ("Biosignatures", "O₂ · O₃ · CH₄ · N₂O", ACCENT2),
    ("Background", "N₂ · H₂", ACCENT3),
    ("Geochemical", "CO · H₂S · SO₂ · NH₃", RGBColor(0x8C, 0x6D, 0x31)),
]
for i, (grp, mols, col) in enumerate(mol_groups):
    y = 4.88 + i*0.44
    add_textbox(sl, grp + ":", 6.6, y, 1.7, 0.38,
                font_size=12, bold=True, color=col)
    add_textbox(sl, mols, 8.35, y, 4.5, 0.38,
                font_size=12, color=DARK_TEXT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — DATA PIPELINE & FEATURE ENGINEERING
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, WHITE)
slide_header(sl, "Data Pipeline & Feature Engineering",
             "Five-stage structured pipeline — each step produces reusable artifacts")
footer(sl)

# Pipeline flow diagram
stages = [
    ("01", "Extract",          "tar.gz → .npy",        NAVY),
    ("02", "Feature Engineer", "Split + Norm + PCA",   ACCENT1),
    ("03", "Baseline RF",      "10k-capped training",  RGBColor(0x5B, 0x8C, 0x5A)),
    ("04", "Deep Model",       "Full-data ResNet",     ACCENT2),
    ("05", "Evaluate",         "Unified comparison",   RGBColor(0x8C, 0x6D, 0x31)),
]

box_w = 2.2
for i, (num, name, desc, col) in enumerate(stages):
    x = 0.25 + i * 2.6
    add_rect(sl, x, 1.35, box_w, 1.5, col)
    add_textbox(sl, num, x+0.08, 1.42, 0.5, 0.45,
                font_size=22, bold=True, color=RGBColor(0xFF, 0xDD, 0x88))
    add_textbox(sl, name, x+0.1, 1.88, box_w-0.2, 0.42,
                font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, desc, x+0.1, 2.3, box_w-0.2, 0.42,
                font_size=11, italic=True, color=RGBColor(0xCC, 0xDD, 0xFF),
                align=PP_ALIGN.CENTER)
    if i < 4:
        add_textbox(sl, "→", x+box_w+0.05, 1.83, 0.4, 0.5,
                    font_size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Feature Engineering detail
add_rect(sl, 0.3, 3.1, 12.7, 3.7, LIGHT_BG)
add_textbox(sl, "Step 2 Deep Dive — Feature Engineering", 0.5, 3.18, 12.0, 0.42,
            font_size=16, bold=True, color=NAVY)

fe_cards = [
    ("70 / 15 / 15 Split", [
        "Deterministic (seed=42)",
        "Fixed once — both models",
        "use identical test set",
        "Saved as train/val/test_indices.npy",
    ], LIGHT_NAVY),
    ("Z-Score Normalisation", [
        "Per CLIMA channel (12 channels)",
        "Fit on train set ONLY",
        "Prevents data leakage",
        "Saved as scaler.joblib",
    ], LIGHT_NAVY),
    ("PCA (RF only)", [
        "Flatten: 12×101 → 1212 features",
        "Reduce to 300 components",
        "Retains ~95% variance",
        "Saved as pca.joblib",
    ], LIGHT_NAVY),
    ("Deep Model Input", [
        "Raw 2D tensor (12 × 101)",
        "No PCA — CNN learns features",
        "Z-normalised spectra only",
        "Stored as spectra_train.npy",
    ], RGBColor(0xE8, 0xF5, 0xE9)),
]

for i, (title, items, fill) in enumerate(fe_cards):
    x = 0.45 + i * 3.15
    add_rect(sl, x, 3.68, 2.9, 2.9, fill)
    add_textbox(sl, title, x+0.12, 3.76, 2.7, 0.38,
                font_size=13, bold=True, color=NAVY)
    for j, item in enumerate(items):
        add_textbox(sl, "· " + item, x+0.15, 4.2 + j*0.46, 2.65, 0.4,
                    font_size=11.5, color=DARK_TEXT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — BASELINE MODEL
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, WHITE)
slide_header(sl, "Baseline Model — Random Forest",
             "12 independent RandomForestRegressors on PCA-reduced features (≤ 10,000 training samples)")
footer(sl)

# Left: how it works
add_rect(sl, 0.3, 1.35, 5.8, 5.6, LIGHT_BG)
add_textbox(sl, "Architecture", 0.5, 1.45, 5.4, 0.38,
            font_size=16, bold=True, color=NAVY)

arch_items = [
    "Input: PCA features (N, 300)",
    "One RandomForestRegressor per molecule",
    "12 completely independent models",
    "Prediction: column-stack 12 scalars → (N, 12)",
]
for i, a in enumerate(arch_items):
    add_textbox(sl, "▸  " + a, 0.45, 1.93 + i*0.52, 5.5, 0.46,
                font_size=13, color=DARK_TEXT)

add_textbox(sl, "Why per-molecule?", 0.5, 4.15, 5.4, 0.38,
            font_size=14, bold=True, color=ACCENT1)
why_items = [
    "Each molecule has different prediction difficulty",
    "Independent tuning of depth / n_estimators",
    "Failing molecule doesn't affect others",
]
for i, w in enumerate(why_items):
    add_textbox(sl, "▸  " + w, 0.45, 4.6 + i*0.45, 5.5, 0.4,
                font_size=12.5, color=DARK_TEXT)

add_textbox(sl, "Why 10k cap?", 0.5, 6.0, 5.4, 0.32,
            font_size=14, bold=True, color=ACCENT2)
add_textbox(sl, "RF is the fast baseline — <1 min training, no GPU needed. Fair comparison point.",
            0.45, 6.38, 5.5, 0.42, font_size=12, italic=True, color=MID_GREY)

# Right: hyperparameter table
add_rect(sl, 6.3, 1.35, 6.7, 5.55, LIGHT_BG)
add_textbox(sl, "Per-Molecule Hyperparameters", 6.5, 1.45, 6.3, 0.38,
            font_size=15, bold=True, color=NAVY)

rf_rows = [
    ("H₂O", "300", "20", "sqrt", "High variance"),
    ("CO₂", "200", "12", "0.5",  "Stable"),
    ("O₂",  "200", "12", "0.5",  "Stable"),
    ("O₃",  "300", "18", "sqrt", "Variable"),
    ("CH₄", "300", "20", "sqrt", "High variance"),
    ("N₂",  "150", "10", "0.4",  "Near-constant"),
    ("N₂O", "300", "18", "sqrt", "Variable"),
    ("CO",  "300", "18", "sqrt", "Variable"),
    ("H₂",  "300", "16", "sqrt", "Variable"),
    ("H₂S", "400", "20", "sqrt", "Trace species"),
    ("SO₂", "400", "22", "sqrt", "Trace species"),
    ("NH₃", "400", "22", "sqrt", "Trace species"),
]

headers_rf = ["Mol", "Trees", "Depth", "max_feat", "Rationale"]
col_w = [0.75, 0.78, 0.72, 0.78, 2.8]
two_col_table(sl, headers_rf, rf_rows,
              left=6.35, top=1.92, width=6.55, height=4.8,
              col_widths=col_w, font_size=11)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — DEEP MODEL ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, WHITE)
slide_header(sl, "Deep Model — 1D ResNet Architecture",
             "Shared convolutional backbone + 12 per-molecule output heads")
footer(sl)

# Architecture flow (vertical blocks)
arch_blocks = [
    ("Input",          "(B, 12, 101)",  "Normalised CLIMA profile — 12 channels × 101 altitude levels", LIGHT_NAVY),
    ("Stem Conv1d",    "(B, 64, 101)",  "Conv(12→64, k=11, s=1, p=5)  +  BatchNorm  +  ReLU", LIGHT_NAVY),
    ("Stage 1  ×2",    "(B,  64, 101)", "ResBlock(64→64, stride=1)  ×2  —  fine-scale altitude features", NAVY),
    ("Stage 2  ×2",    "(B, 128,  51)", "ResBlock(64→128, stride=2)  ×2  —  halve sequence length", NAVY),
    ("Stage 3  ×2",    "(B, 256,  26)", "ResBlock(128→256, stride=2)  ×2", ACCENT1),
    ("Stage 4  ×2",    "(B, 512,  13)", "ResBlock(256→512, stride=2)  ×2  —  coarse global features", ACCENT1),
    ("GlobalAvgPool",  "(B, 512)",      "AdaptiveAvgPool1d(1) — squeeze altitude dimension", ACCENT2),
    ("Shared FC",      "(B, 256)",      "Dropout(0.3)  +  FC(512→256)  +  LayerNorm  +  ReLU", ACCENT2),
    ("12 × Head",      "(B, 12)",       "Per-molecule MLP: FC→LN→ReLU→Dropout → scalar log₁₀ abundance", RGBColor(0x1A, 0x7A, 0x3C)),
]

block_h = 0.57
for i, (name, shape, desc, col) in enumerate(arch_blocks):
    y = 1.35 + i * block_h
    add_rect(sl, 0.3,  y, 2.35, block_h-0.04, col)
    add_rect(sl, 2.7,  y, 1.55, block_h-0.04, LIGHT_BG)
    add_rect(sl, 4.3,  y, 8.7,  block_h-0.04, LIGHT_BG)

    add_textbox(sl, name,  0.38, y+0.06, 2.2, block_h-0.12,
                font_size=12, bold=True, color=WHITE if col != LIGHT_NAVY else NAVY)
    add_textbox(sl, shape, 2.75, y+0.06, 1.5, block_h-0.12,
                font_size=11, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)
    add_textbox(sl, desc,  4.38, y+0.06, 8.5, block_h-0.12,
                font_size=12, color=DARK_TEXT)

# Col headers
add_rect(sl, 0.3, 1.28, 2.35, 0.07, RED)
add_rect(sl, 2.7, 1.28, 1.55, 0.07, RED)
add_rect(sl, 4.3, 1.28, 8.7,  0.07, RED)
add_textbox(sl, "Layer", 0.38, 1.18, 2.2, 0.3, font_size=11, bold=True, color=NAVY)
add_textbox(sl, "Output shape", 2.75, 1.18, 1.5, 0.3, font_size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_textbox(sl, "Description", 4.38, 1.18, 8.5, 0.3, font_size=11, bold=True, color=NAVY)

# ResBlock detail in margin — right side annotation
add_textbox(sl, "ResBlock1D:", 10.8, 1.35, 2.3, 0.3, font_size=10, bold=True, color=ACCENT1)
resblock_desc = "Conv1d→BN→ReLU\n Conv1d→BN\n+ skip (x) → ReLU\n Shortcut = 1×1 Conv\n if size changes"
add_textbox(sl, resblock_desc, 10.8, 1.65, 2.3, 1.8,
            font_size=10, italic=True, color=MID_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — DEEP MODEL: KEY DESIGN CHOICES
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, WHITE)
slide_header(sl, "Deep Model — Key Design Choices",
             "Every design decision is motivated by the data's properties")
footer(sl)

choices = [
    ("1D Convolutions", ACCENT1,
     "CLIMA profiles are 1D sequences over altitude. 1D CNNs detect local patterns "
     "(e.g. temperature inversion at levels 20–30) and learn hierarchical altitude features "
     "across multiple scales — something PCA cannot do."),
    ("Residual Skip Connections", NAVY,
     "Prevents vanishing gradients in deep networks. Output = F(x) + x — the network learns "
     "the residual correction rather than the full transformation. Enables stable training "
     "of 9-layer+ networks without gradient decay."),
    ("Shared Backbone", ACCENT2,
     "The ResNet body learns general atmospheric representations that are useful for all 12 "
     "molecules simultaneously — multi-task learning. The shared 256-d embedding carries "
     "information about atmospheric state that each head then specialises."),
    ("Per-Molecule Heads", RGBColor(0x1A, 0x7A, 0x3C),
     "Each molecule has its own MLP head with tuned depth and dropout. Trace species (H₂S, "
     "SO₂, NH₃) get 3-layer heads with higher dropout (0.40) to prevent overfitting. Stable "
     "species (CO₂, O₂) get lighter 2-layer heads with lower dropout (0.20)."),
    ("Weighted MSE Loss", RGBColor(0x8C, 0x6D, 0x31),
     "Standard MSE would let easy molecules dominate. Weights (1.0–2.0) upweight trace "
     "species (SO₂, NH₃ at 2.0) and downweight near-constant N₂ (1.0), forcing the "
     "network to invest in the molecules that are hardest to detect."),
    ("Gaussian Noise Augmentation", RGBColor(0x55, 0x55, 0x55),
     "Adding noise std=0.01 to training spectra prevents memorisation of exact profile "
     "values. Forces the ResNet to learn robust features that generalise to unseen "
     "atmospheric states — analogous to image augmentation in vision CNNs."),
]

for i, (title, col, body) in enumerate(choices):
    row = i // 2
    col_i = i % 2
    x = 0.3 + col_i * 6.5
    y = 1.35 + row * 1.95
    add_rect(sl, x, y, 6.2, 1.82, LIGHT_BG)
    add_rect(sl, x, y, 0.12, 1.82, col)
    add_textbox(sl, title, x+0.22, y+0.12, 5.8, 0.38,
                font_size=14, bold=True, color=col)
    add_textbox(sl, body, x+0.22, y+0.54, 5.8, 1.18,
                font_size=11.5, color=DARK_TEXT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — TRAINING STRATEGY
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, WHITE)
slide_header(sl, "Training Strategy",
             "AdamW · Cosine LR schedule · Early stopping · Gradient clipping")
footer(sl)

# Left: training config
add_rect(sl, 0.3, 1.35, 5.7, 5.55, LIGHT_BG)
add_textbox(sl, "Training Configuration", 0.5, 1.45, 5.3, 0.38,
            font_size=16, bold=True, color=NAVY)

config_rows = [
    ("Optimizer",         "AdamW"),
    ("Learning rate",     "1 × 10⁻³"),
    ("Weight decay",      "1 × 10⁻⁴"),
    ("Scheduler",         "CosineAnnealingLR"),
    ("T_max / eta_min",   "150 / 1 × 10⁻⁶"),
    ("Max epochs",        "150"),
    ("Early stop patience", "30 epochs"),
    ("Min delta",         "1 × 10⁻⁵"),
    ("Batch size",        "32 (train) / 64 (eval)"),
    ("Grad clip norm",    "5.0"),
    ("Noise augmentation","std = 0.01"),
    ("Dataset (train)",   "~87,000 samples"),
]
for i, (k, v) in enumerate(config_rows):
    y = 1.95 + i*0.39
    add_textbox(sl, k, 0.45, y, 3.0, 0.35, font_size=12.5, bold=True, color=NAVY)
    add_textbox(sl, v, 3.5,  y, 2.3, 0.35, font_size=12.5, color=DARK_TEXT)

# Right: explanations
explanations = [
    ("AdamW", ACCENT1,
     "Adam with decoupled weight decay — L2 regularisation is applied directly to weights, "
     "independent of the adaptive gradient. Better generalisation than standard Adam."),
    ("Cosine Annealing", NAVY,
     "LR starts at 1e-3, smoothly decays to 1e-6 following a cosine curve. "
     "No sharp drops. Naturally slows near minima to avoid oscillation."),
    ("Early Stopping", ACCENT2,
     "Monitors val_loss. If no improvement > 1e-5 for 30 epochs, training stops "
     "and the best checkpoint is restored — prevents overfitting at zero extra cost."),
    ("Gradient Clipping", RGBColor(0x1A, 0x7A, 0x3C),
     "If gradient norm > 5.0, all gradients are scaled down proportionally. "
     "Prevents exploding gradient spikes that destabilise early training."),
]

for i, (title, col, body) in enumerate(explanations):
    y = 1.35 + i * 1.38
    add_rect(sl, 6.3, y, 6.7, 1.28, LIGHT_BG)
    add_rect(sl, 6.3, y, 0.1, 1.28, col)
    add_textbox(sl, title, 6.5, y+0.1, 6.3, 0.36, font_size=14, bold=True, color=col)
    add_textbox(sl, body, 6.5, y+0.5, 6.3, 0.72, font_size=12, color=DARK_TEXT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — EVALUATION METHODOLOGY
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, WHITE)
slide_header(sl, "Evaluation Methodology",
             "Held-out test set · Per-molecule R², RMSE, MAE · Unified comparison report")
footer(sl)

# Metric cards
metrics = [
    ("R²\nCoefficient of Determination",
     "1 − Σ(y − ŷ)² / Σ(y − ȳ)²",
     "Primary metric. Scale-normalised — 0.75 means '75% of variance explained' "
     "regardless of molecule. Enables fair cross-molecule comparison. R² < 0 means "
     "worse than predicting the mean.",
     ACCENT1),
    ("RMSE\nRoot Mean Squared Error",
     "√[Σ(y − ŷ)² / N]",
     "In log₁₀ units (orders of magnitude). Penalises large errors more than MAE. "
     "A RMSE of 0.5 means predictions are typically ½ an order of magnitude off.",
     NAVY),
    ("MAE\nMean Absolute Error",
     "Σ|y − ŷ| / N",
     "Also in log₁₀ units. More robust to outliers than RMSE. Complementary to RMSE — "
     "if RMSE >> MAE, there are a few large prediction errors dominating.",
     ACCENT2),
]

for i, (title, formula, desc, col) in enumerate(metrics):
    x = 0.3 + i * 4.35
    add_rect(sl, x, 1.35, 4.1, 2.75, col)
    add_textbox(sl, title, x+0.15, 1.43, 3.8, 0.72,
                font_size=14, bold=True, color=WHITE)
    add_rect(sl, x+0.15, 2.2, 3.8, 0.48, RGBColor(0xFF, 0xFF, 0xFF))
    add_textbox(sl, formula, x+0.2, 2.24, 3.7, 0.42,
                font_size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_textbox(sl, desc, x+0.15, 2.76, 3.8, 1.22, font_size=12, color=WHITE)

# Fair comparison principles
add_rect(sl, 0.3, 4.3, 12.7, 2.45, LIGHT_BG)
add_textbox(sl, "Evaluation Fairness Principles", 0.5, 4.38, 12.0, 0.38,
            font_size=16, bold=True, color=NAVY)

principles = [
    ("Same test set", "Both models evaluated on the identical 19k-sample test split, fixed in Step 2 with seed=42. Scores are directly comparable."),
    ("Intentionally asymmetric training", "RF uses 10k samples; ResNet uses ~87k. The comparison is baseline vs deep model — not equal-data ablation."),
    ("Per-molecule granularity", "Reporting R²/RMSE/MAE for each of 12 molecules reveals which species each model handles well or poorly."),
    ("Mean R² summary", "Unweighted mean across 12 molecules provides a single-number comparison. NH₃ and N₂ drag it down for both models."),
]
for i, (k, v) in enumerate(principles):
    col_i = i % 2
    row_i = i // 2
    x = 0.45 + col_i * 6.3
    y = 4.85 + row_i * 0.68
    add_textbox(sl, k + ": ", x, y, 2.1, 0.35, font_size=12.5, bold=True, color=ACCENT1)
    add_textbox(sl, v, x+2.15, y, 3.95, 0.62, font_size=12, color=DARK_TEXT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — RESULTS
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, WHITE)
slide_header(sl, "Results — Per-Molecule Test R²",
             "Test set: ~19,000 samples — identical for both models")
footer(sl)

# Results table
results_data = [
    ("H₂O",  "+0.3669", "+0.6299", "+0.26", "Deep wins — ResNet captures altitude patterns"),
    ("CO₂",  "+0.3262", "+0.6006", "+0.27", "Deep wins — non-linear vertical CO₂ structure"),
    ("O₂",   "+0.5985", "+0.5939", "−0.00", "Tie — linear PCA sufficient for O₂"),
    ("O₃",   "+0.7547", "+0.8954", "+0.14", "Best molecule for both — strong signal"),
    ("CH₄",  "+0.4793", "+0.6978", "+0.22", "Deep wins — high variance, complex patterns"),
    ("N₂",   "−0.0421", "+0.0223", "+0.06", "Both fail — near-constant, no signal"),
    ("N₂O",  "+0.2133", "+0.1137", "−0.10", "RF wins — deep model overfits low-variance signal"),
    ("CO",   "+0.0465", "+0.0006", "−0.05", "Both fail — complex photochemistry not in features"),
    ("H₂",   "+0.2008", "+0.0187", "−0.18", "RF wins — deep model overfits"),
    ("H₂S",  "+0.5461", "+0.8162", "+0.27", "Deep wins — distinctive altitude signature"),
    ("SO₂",  "+0.2198", "+0.3538", "+0.13", "Deep wins — trace species with clear signal"),
    ("NH₃",  "+0.4107", "+0.5652", "+0.16", "Both predict floor (−40) — not truly learnable"),
    ("MEAN", "+0.3434", "+0.4423", "+0.10", "Overall: +9.9% mean R² improvement"),
]

headers_r = ["Molecule", "Baseline R²", "ResNet R²", "ΔR²", "Interpretation"]
col_w_r = [1.1, 1.3, 1.2, 1.0, 5.4]
row_h = 0.35

# Draw table
y_start = 1.35
x_start = 0.3

# Header row
x = x_start
for h, cw in zip(headers_r, col_w_r):
    add_rect(sl, x, y_start, cw, row_h, NAVY)
    add_textbox(sl, h, x+0.04, y_start+0.04, cw-0.08, row_h-0.06,
                font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    x += cw

for ri, row in enumerate(results_data):
    y = y_start + row_h * (ri + 1)
    is_mean = row[0] == "MEAN"
    delta_str = row[3].replace("−", "-").replace("+", "")
    delta = float(delta_str) if delta_str not in ("", "—") else 0.0
    x = x_start

    for ci, (val, cw) in enumerate(zip(row, col_w_r)):
        if is_mean:
            fill = NAVY
            txt_col = WHITE
        elif ci == 1:  # baseline
            fill = LIGHT_BG
            txt_col = DARK_TEXT
        elif ci == 2:  # deep
            r2 = float(val)
            fill = RGBColor(0xE8, 0xF5, 0xE9) if r2 > 0.5 else (
                   RGBColor(0xFF, 0xF8, 0xE1) if r2 > 0.2 else
                   RGBColor(0xFF, 0xEE, 0xEE))
            txt_col = DARK_TEXT
        elif ci == 3:  # delta
            fill = RGBColor(0xE8, 0xF5, 0xE9) if delta > 0 else RGBColor(0xFF, 0xEE, 0xEE)
            txt_col = ACCENT3 if delta > 0 else ACCENT2
        else:
            fill = LIGHT_BG if ri % 2 == 0 else WHITE
            txt_col = DARK_TEXT if not is_mean else WHITE

        add_rect(sl, x, y, cw, row_h, fill)
        add_textbox(sl, val, x+0.04, y+0.04, cw-0.08, row_h-0.06,
                    font_size=10.5 if ci == 4 else 11,
                    bold=(ci == 0 or is_mean),
                    color=WHITE if is_mean else txt_col,
                    align=PP_ALIGN.LEFT if ci == 4 else PP_ALIGN.CENTER)
        x += cw


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — KEY FINDINGS
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, WHITE)
slide_header(sl, "Key Findings & Insights",
             "What the results tell us about ML-based atmospheric retrieval")
footer(sl)

findings = [
    ("🌟", "Deep model wins on 8/12 molecules", ACCENT2,
     "The 1D ResNet outperforms the RF baseline on most molecules, demonstrating that "
     "convolutional feature learning from altitude profiles captures patterns that PCA discards."),
    ("📈", "Largest gains on complex molecules", ACCENT1,
     "H₂O (+26%), CO₂ (+27%), H₂S (+27%), CH₄ (+22%) — molecules with rich altitude-dependent "
     "chemistry benefit most from the spatial feature extraction of the ResNet."),
    ("⚖️", "RF remains competitive on stable species", RGBColor(0x1A, 0x7A, 0x3C),
     "O₂ is effectively tied (ΔR²=−0.005). For near-constant species, the linear PCA projection "
     "captures the dominant variance — deeper architectures don't help here."),
    ("⚠️", "Both models fail on CO, H₂, N₂", RGBColor(0x8C, 0x6D, 0x31),
     "These molecules have R²≈0. Their abundance depends on fine-grained photochemical balance "
     "not captured in CLIMA profiles — a fundamental input feature limitation, not a model issue."),
    ("🔬", "Biosignatures are learnable", NAVY,
     "O₂ (0.60), O₃ (0.90), CH₄ (0.70), N₂O (0.11). O₂ and O₃ are well-predicted. "
     "CH₄ shows strong improvement. N₂O remains difficult — future work needed."),
    ("💡", "Multi-task learning helps", RGBColor(0x55, 0x55, 0x55),
     "The shared backbone implicitly benefits from being trained on all 12 targets "
     "simultaneously. Cross-molecule correlations (e.g. O₂/O₃ co-variation) are captured "
     "in the shared representation."),
]

for i, (icon, title, col, body) in enumerate(findings):
    row = i // 2
    ci = i % 2
    x = 0.3 + ci * 6.5
    y = 1.35 + row * 1.85
    add_rect(sl, x, y, 6.2, 1.72, LIGHT_BG)
    add_rect(sl, x, y, 0.12, 1.72, col)
    add_textbox(sl, icon + "  " + title, x+0.22, y+0.1, 5.8, 0.4,
                font_size=13, bold=True, color=col)
    add_textbox(sl, body, x+0.22, y+0.55, 5.8, 1.05, font_size=12, color=DARK_TEXT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — MODEL COMPARISON SUMMARY
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, WHITE)
slide_header(sl, "Model Comparison Summary",
             "Baseline RF vs SpectralResNet — design, compute, and performance")
footer(sl)

# Comparison table
headers_c = ["Aspect", "Random Forest (Baseline)", "SpectralResNet (Deep)"]
rows_c = [
    ("Architecture",       "12 independent RFs",              "1D ResNet + 12 heads"),
    ("Input format",       "Flat PCA vector (300-d)",          "2D tensor (12 × 101)"),
    ("Feature extraction", "PCA — linear, pre-computed",       "Learned convolutional filters"),
    ("Multi-task",         "None — fully independent",         "Yes — shared backbone"),
    ("Training samples",   "10,000 (hard cap)",                "~87,000 (full training split)"),
    ("Training time",      "< 1 min (CPU)",                    "10–20 min (MPS/GPU)"),
    ("Parameters",         "Non-parametric",                   "~2 million"),
    ("GPU required",       "No",                               "Recommended"),
    ("Interpretability",   "Feature importance available",     "Black box (lower)"),
    ("Mean test R²",       "0.34",                             "0.44  (+29% relative)"),
    ("Best molecule",      "O₃  R²=0.75",                      "O₃  R²=0.90"),
    ("Worst molecule",     "N₂  R²=−0.04",                     "CO  R²=0.00"),
]

col_w_c = [3.0, 4.8, 4.8]
two_col_table(sl, headers_c, rows_c,
              left=0.3, top=1.35, width=12.73, height=5.5,
              col_widths=col_w_c, font_size=12)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — CHALLENGES & LIMITATIONS
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, WHITE)
slide_header(sl, "Challenges & Limitations",
             "Honest assessment of what works and what doesn't — and why")
footer(sl)

challenges = [
    ("Data scale vs compute", ACCENT2,
     "Full INARA has 3.1M spectra. We used ~124k due to local compute constraints (~10 GB "
     "processed). Extraction from 10 tar.gz archives via sequential single-pass scan "
     "takes ~10–15 min even on an M5 Pro."),
    ("NH₃ — not learnable from CLIMA", NAVY,
     "NH₃ values are at the log-floor (−40) for nearly all samples in this dataset. "
     "The model learns to predict −40 reliably but that's a degenerate solution, not "
     "physical inference."),
    ("CO and H₂ — fundamental degeneracy", ACCENT1,
     "These molecules' abundances depend on photochemical fine structure not captured "
     "in CLIMA profiles. No ML model can overcome missing input information — this is an "
     "inverse problem degeneracy, not a model failure."),
    ("Asymmetric comparison", RGBColor(0x1A, 0x7A, 0x3C),
     "The RF uses 10k training samples; the ResNet uses 87k. The +9.9% R² improvement "
     "is partly attributable to more data, not purely to the architecture. A data-equal "
     "ablation would isolate the architecture contribution."),
    ("N₂ near-constant challenge", RGBColor(0x8C, 0x6D, 0x31),
     "N₂ std = 0.044 in log₁₀ space — essentially constant. R² is a poor metric here "
     "because the denominator (variance) is near-zero. RMSE shows the model does "
     "reasonably, but R² is misleading."),
    ("Real data validation pending", RGBColor(0x55, 0x55, 0x55),
     "All experiments are on synthetic INARA spectra. Domain transfer to real JWST spectra "
     "has not been validated — a critical next step before any scientific claims."),
]

for i, (title, col, body) in enumerate(challenges):
    row = i // 2
    ci = i % 2
    x = 0.3 + ci * 6.5
    y = 1.35 + row * 1.88
    add_rect(sl, x, y, 6.2, 1.76, LIGHT_BG)
    add_rect(sl, x, y, 0.12, 1.76, col)
    add_textbox(sl, title, x+0.22, y+0.1, 5.8, 0.38,
                font_size=13, bold=True, color=col)
    add_textbox(sl, body, x+0.22, y+0.52, 5.8, 1.12, font_size=12, color=DARK_TEXT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — FUTURE WORK
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, WHITE)
slide_header(sl, "Future Work",
             "Extending this project toward real-world scientific applicability")
footer(sl)

future = [
    ("🔭", "JWST Domain Transfer", ACCENT1,
     "Test trained models on real JWST transmission spectra and compare predictions "
     "to published Bayesian retrievals. Key validation step for scientific credibility."),
    ("📊", "Full 3.1M Dataset", NAVY,
     "Train the ResNet on the complete INARA dataset using Northeastern Explorer HPC. "
     "Expected improvement in R² for high-variance molecules like H₂O and CH₄."),
    ("🧠", "Uncertainty Quantification", ACCENT2,
     "Add Monte Carlo Dropout or deep ensembles to provide calibrated confidence intervals "
     "— critical for scientific use (e.g. 'O₂ = −0.7 ± 0.2 log₁₀')."),
    ("🔬", "Attention / Transformer", RGBColor(0x1A, 0x7A, 0x3C),
     "Replace or augment the ResNet with a 1D Transformer to learn long-range "
     "altitude-level dependencies. Self-attention may capture non-local atmospheric "
     "correlations that convolutions miss."),
    ("⚖️", "Equal-Data Ablation", RGBColor(0x8C, 0x6D, 0x31),
     "Train the RF baseline on 87k samples (same as ResNet) to cleanly isolate the "
     "architecture benefit from the data quantity benefit."),
    ("🌍", "Four Stellar Types", RGBColor(0x55, 0x55, 0x55),
     "Stratify analysis by stellar type (F, G, K, M). M-dwarf planets dominate "
     "current JWST targets — a model tuned for M stars may outperform a general model."),
]

for i, (icon, title, col, body) in enumerate(future):
    row = i // 2
    ci = i % 2
    x = 0.3 + ci * 6.5
    y = 1.35 + row * 1.88
    add_rect(sl, x, y, 6.2, 1.76, LIGHT_BG)
    add_rect(sl, x, y, 0.12, 1.76, col)
    add_textbox(sl, icon + "  " + title, x+0.22, y+0.1, 5.8, 0.38,
                font_size=13, bold=True, color=col)
    add_textbox(sl, body, x+0.22, y+0.52, 5.8, 1.12, font_size=12, color=DARK_TEXT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — CONCLUSION
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, NAVY)
add_rect(sl, 0, 0, 0.18, 7.5, RED)

add_textbox(sl, "Conclusion", 0.5, 0.4, 12.5, 0.65,
            font_size=34, bold=True, color=WHITE)
add_textbox(sl, "CS 6140 · Machine Learning · Northeastern University · Spring 2026",
            0.5, 1.05, 12.5, 0.38, font_size=14, italic=True,
            color=RGBColor(0xAA, 0xBB, 0xDD))

# Red separator
add_rect(sl, 0.5, 1.52, 12.0, 0.06, RED)

conclusions = [
    ("We built an end-to-end ML pipeline", "for exoplanet atmospheric retrieval — from raw synthetic spectra to per-molecule R² metrics."),
    ("Two models were compared:", "Random Forest baseline (10k cap) and 1D SpectralResNet (full training split)."),
    ("ResNet achieves mean R² = 0.44", "vs baseline 0.34 — a +29% relative improvement — using the same held-out test set."),
    ("Largest gains on biosignatures:", "O₃ (0.90), CH₄ (0.70), H₂O (0.63), O₂ (0.59) — the molecules most relevant to the search for life."),
    ("ML retrieval is viable and fast:", "milliseconds per planet vs. hours for Bayesian methods — enabling population-level studies."),
    ("Next step — real JWST data:", "validate on actual observations and add uncertainty quantification for scientific deployment."),
]

for i, (bold_part, rest) in enumerate(conclusions):
    y = 1.72 + i * 0.75
    add_rect(sl, 0.5, y, 0.5, 0.58, RED)
    add_textbox(sl, str(i+1), 0.52, y+0.06, 0.45, 0.45,
                font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, bold_part + "  " + rest, 1.15, y+0.06,
                11.5, 0.5, font_size=14.5, color=WHITE)

# Team names at bottom
add_rect(sl, 0.5, 7.0, 12.0, 0.35, RGBColor(0x00, 0x1A, 0x3A))
add_textbox(sl,
    "Shantanu Wankhare  ·  Bhalchandra Shinde  ·  Asad Mulani   |   Dataset: INARA (NASA FDL / Zorzan et al. 2025)",
    0.55, 7.02, 11.8, 0.3, font_size=10.5, color=RGBColor(0xAA, 0xBB, 0xDD),
    align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — REFERENCES
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, WHITE)
slide_header(sl, "References")
footer(sl)

refs = [
    ("[1]", "Zorzan et al. (2025)", "INARA: A Machine Learning Dataset for Exoplanet Atmospheric Retrieval. NASA FDL / Exoplanet Archive.", "Dataset"),
    ("[2]", "He et al. (2016)", "Deep Residual Learning for Image Recognition. CVPR 2016.", "ResNet architecture"),
    ("[3]", "Loshchilov & Hutter (2019)", "Decoupled Weight Decay Regularisation. ICLR 2019.", "AdamW optimizer"),
    ("[4]", "Breiman (2001)", "Random Forests. Machine Learning, 45(1), 5–32.", "Random Forest baseline"),
    ("[5]", "Márquez-Neila et al. (2018)", "Supervised machine learning for analysing spectra of exoplanetary atmospheres. Nature Astronomy.", "ML retrieval"),
    ("[6]", "Vasist et al. (2023)", "Neural posterior estimation for exoplanetary atmospheric retrieval. A&A, 672.", "Deep learning retrieval"),
    ("[7]", "Gebhard et al. (2024)", "Simulation-based inference for exoplanet atmospheric retrievals. A&A.", "SBI retrieval"),
    ("[8]", "JWST Collaboration (2023)", "Early Release Science observations of transiting exoplanet atmospheres.", "Real spectra context"),
]

for i, (num, authors, title, tag) in enumerate(refs):
    y = 1.45 + i * 0.67
    add_rect(sl, 0.3, y, 0.55, 0.52, NAVY)
    add_textbox(sl, num, 0.32, y+0.06, 0.5, 0.4,
                font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, authors, 0.95, y+0.02, 2.8, 0.25, font_size=12, bold=True, color=NAVY)
    add_textbox(sl, title, 0.95, y+0.26, 10.5, 0.22, font_size=11, color=DARK_TEXT)
    add_rect(sl, 11.5, y, 1.5, 0.52, LIGHT_BG)
    add_textbox(sl, tag, 11.55, y+0.08, 1.4, 0.38, font_size=10, italic=True, color=MID_GREY)


# ── Save ──────────────────────────────────────────────────────────────────────
OUT_PATH = Path(__file__).parent / 'INARA_Presentation.pptx'
prs.save(str(OUT_PATH))
print(f'Saved: {OUT_PATH}')
print(f'Slides: {len(prs.slides)}')
